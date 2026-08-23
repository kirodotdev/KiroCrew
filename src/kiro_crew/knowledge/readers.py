"""File readers for knowledge ingestion. Supports text, PDF, PPTX, DOCX, HTML."""

import codecs
import os
import re
from pathlib import Path

from kiro_crew.security import is_sensitive_path

try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

try:
    from docx import Document  # type: ignore[import-untyped]
except ImportError:
    Document = None  # type: ignore[assignment,misc]

try:
    import html2text as _html2text_mod
except ImportError:
    _html2text_mod = None  # type: ignore[assignment]


def _decode_text_bytes(raw: bytes) -> str:
    """Decode raw file bytes as text.

    BOM-sniffed UTF-16 LE/BE first: Windows PowerShell tooling
    (New-ModuleManifest, the legacy ISE) writes UTF-16LE, whose bytes miss
    utf-8 and would land in the latin-1 fallback as BOM + interleaved NULs --
    mojibake in the index. The branch is extension-agnostic: any text format
    arriving as BOM'd UTF-16 decodes correctly. A BOM that lies (truncated
    copies, binary that happens to start FF FE) degrades to latin-1 rather
    than failing the file's ingest. Everything else keeps the historical
    utf-8 -> latin-1 chain. Operating on one in-memory buffer (never
    reopening the path) keeps the sensitive-path check race-free, and the
    newline translation mirrors text-mode open() so results are unchanged.
    """
    if raw[:2] in (codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE):
        try:
            text = raw.decode('utf-16')
        except UnicodeDecodeError:
            text = raw.decode('latin-1')
    else:
        try:
            text = raw.decode('utf-8')
        except UnicodeDecodeError:
            text = raw.decode('latin-1')
    return text.replace('\r\n', '\n').replace('\r', '\n')


def _read_error(exc: Exception) -> tuple[str, dict]:
    """The reader result for a file that could not be read.

    ``format: 'error'`` is the sentinel ``IngestionPipeline.ingest_file`` tests:
    it raises rather than indexing, so the file is recorded as failed against its
    source and the surrounding scan continues. The returned text is diagnostic and
    never becomes an item.
    """
    return f'Error reading file: {exc}', {'format': 'error', 'error': str(exc)}


def _missing_dep(fmt: str, package: str) -> tuple[str, dict]:
    """The reader result when an optional extraction dependency is absent.

    Carries the same ``format: 'error'`` sentinel as :func:`_read_error`. ``.pptx``
    has no declared dependency, so this is the normal path for that format rather
    than an anomaly.
    """
    message = f'{fmt} support requires {package}: pip install {package}'
    return message, {'format': 'error', 'error': f'{fmt} support requires {package}'}


class FileReader:
    # Binary formats need optional runtime deps: .pdf -> pdfplumber and .docx ->
    # python-docx (both declared in setup.cfg). .pptx -> python-pptx is NOT declared,
    # so .pptx is intentionally kept out of SUPPORTED even though _read_pptx exists.
    SUPPORTED = {
        '', '.md', '.txt', '.org', '.py', '.java', '.ts', '.js', '.rs', '.go',
        '.html', '.htm', '.docx', '.pdf',
        '.csv', '.log', '.json', '.jsonl', '.ndjson', '.yaml', '.yml',
        '.sh', '.rb', '.ps1', '.psm1', '.psd1', '.c', '.cpp', '.h',
    }

    _DISPATCH = {
        '.pdf': '_read_pdf',
        '.pptx': '_read_pptx',
        '.docx': '_read_docx',
        '.html': '_read_html',
        '.htm': '_read_html',
    }

    def read(self, path: str) -> tuple[str, dict]:
        if is_sensitive_path(path):
            raise PermissionError(f"Refusing to read sensitive path: {path}")
        p = Path(path)
        ext = p.suffix.lower()
        base_meta = {
            'format': ext.lstrip('.'),
            'title': p.stem,
            'file_size': os.path.getsize(path),
            'extension': ext,
        }
        method_name = self._DISPATCH.get(ext)
        if method_name:
            text, meta = getattr(self, method_name)(path)
        else:
            text, meta = self._read_text(path, ext.lstrip('.'))
        base_meta.update(meta)
        base_meta['line_count'] = text.count('\n') + 1 if text else 0
        return text, base_meta

    def _read_text(self, path: str, fmt: str) -> tuple[str, dict]:
        try:
            # One open, one read: decoding the buffer (rather than reopening
            # the path per encoding attempt) means the file that was
            # sensitive-path-checked is the file that gets indexed.
            with open(path, 'rb') as f:
                raw = f.read()
            return _decode_text_bytes(raw), {'format': fmt}
        except Exception as e:
            return _read_error(e)

    def _read_pdf(self, path: str) -> tuple[str, dict]:
        if pdfplumber is None:
            return _missing_dep('PDF', 'pdfplumber')
        try:
            with pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or '' for p in pdf.pages]
                return '\n'.join(pages), {'format': 'pdf', 'page_count': len(pages)}
        except Exception as e:
            return _read_error(e)

    def _read_pptx(self, path: str) -> tuple[str, dict]:
        if Presentation is None:
            return _missing_dep('PPTX', 'python-pptx')
        try:
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                title = ''
                body_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if shape == slide.shapes.title:
                            title = text
                        else:
                            body_parts.append(text)
                notes = ''
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                section = f'## Slide {i}: {title}\n{chr(10).join(body_parts)}'
                if notes:
                    section += f'\n{notes}'
                parts.append(section)
            return '\n\n'.join(parts), {'format': 'pptx', 'slide_count': len(prs.slides)}
        except Exception as e:
            return _read_error(e)

    def _read_docx(self, path: str) -> tuple[str, dict]:
        if Document is None:
            return _missing_dep('DOCX', 'python-docx')
        try:
            doc = Document(path)
            lines = []
            for para in doc.paragraphs:
                style = para.style.name if para.style else ''
                text = para.text
                if style.startswith('Heading'):
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    lines.append(f'{"#" * level} {text}')
                else:
                    lines.append(text)
            return '\n'.join(lines), {'format': 'docx', 'content_type': 'markdown', 'paragraph_count': len(doc.paragraphs)}
        except Exception as e:
            return _read_error(e)

    def _read_html(self, path: str) -> tuple[str, dict]:
        try:
            # Same single-open buffer decode as _read_text: HTML saved by
            # Windows tooling can arrive as BOM'd UTF-16 too.
            with open(path, 'rb') as f:
                html = _decode_text_bytes(f.read())
        except Exception as e:
            return _read_error(e)
        if _html2text_mod is not None:
            h = _html2text_mod.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            return h.handle(html), {'format': 'html'}
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text, {'format': 'html'}
